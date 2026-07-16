import csv
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


@dataclass
class Block:
    text: str
    page: Optional[int] = None
    heading: Optional[str] = None


@dataclass
class ParsedDocument:
    blocks: list[Block] = field(default_factory=list)
    title: Optional[str] = None
    author: Optional[str] = None
    page_count: int = 0


class UnsupportedFileType(Exception):
    pass


def parse_file(path: str) -> ParsedDocument:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    if ext == ".csv":
        return _parse_csv(path)
    if ext in (".txt", ".md"):
        return _parse_text(path)
    if ext == ".rtf":
        return _parse_rtf(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext in (".html", ".htm"):
        return _parse_html(path)
    raise UnsupportedFileType(ext)


def _parse_pdf(path: str) -> ParsedDocument:
    import fitz

    doc = fitz.open(path)
    blocks: list[Block] = []
    title = doc.metadata.get("title") or None
    author = doc.metadata.get("author") or None
    for page_index in range(len(doc)):
        page = doc[page_index]
        page_no = page_index + 1
        text = page.get_text("text").strip()
        if not text:
            text = _ocr_pdf_page(page)
        if not text:
            continue
        heading = _guess_heading(text)
        blocks.append(Block(text=text, page=page_no, heading=heading))
    page_count = len(doc)
    doc.close()
    return ParsedDocument(blocks=blocks, title=title, author=author, page_count=page_count)


def _ocr_pdf_page(page) -> str:
    try:
        import io

        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=200)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ""


def _parse_docx(path: str) -> ParsedDocument:
    import docx

    document = docx.Document(path)
    blocks: list[Block] = []
    current_heading: Optional[str] = None
    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading") or style == "title":
            current_heading = text
        blocks.append(Block(text=text, heading=current_heading))
    for ti, table in enumerate(document.tables):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            blocks.append(
                Block(text="\n".join(rows), heading=f"Table {ti + 1}")
            )
    props = document.core_properties
    title = props.title or None
    author = props.author or None
    return ParsedDocument(blocks=blocks, title=title, author=author)


def _parse_xlsx(path: str) -> ParsedDocument:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    blocks: list[Block] = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(Block(text="\n".join(rows), heading=f"Sheet: {sheet.title}"))
    wb.close()
    return ParsedDocument(blocks=blocks)


def _parse_csv(path: str) -> ParsedDocument:
    rows = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            cells = [c.strip() for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
    text = "\n".join(rows)
    return ParsedDocument(blocks=[Block(text=text)] if text else [])


def _parse_text(path: str) -> ParsedDocument:
    with open(path, encoding="utf-8", errors="ignore") as f:
        content = f.read()
    blocks: list[Block] = []
    current_heading: Optional[str] = None
    for para in content.split("\n\n"):
        text = para.strip()
        if not text:
            continue
        if text.startswith("#"):
            current_heading = text.splitlines()[0].lstrip("#").strip()
        blocks.append(Block(text=text, heading=current_heading))
    return ParsedDocument(blocks=blocks)


def _parse_rtf(path: str) -> ParsedDocument:
    from striprtf.striprtf import rtf_to_text

    with open(path, encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    text = rtf_to_text(raw)
    blocks = [Block(text=p.strip()) for p in text.split("\n\n") if p.strip()]
    return ParsedDocument(blocks=blocks)


def _parse_pptx(path: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(path)
    blocks: list[Block] = []
    for idx, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            blocks.append(
                Block(text="\n".join(texts), page=idx + 1, heading=f"Slide {idx + 1}")
            )
    return ParsedDocument(blocks=blocks, page_count=len(prs.slides))


def _parse_html(path: str) -> ParsedDocument:
    from bs4 import BeautifulSoup

    with open(path, encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    for tag in soup(["script", "style"]):
        tag.decompose()
    title = soup.title.string.strip() if soup.title and soup.title.string else None
    blocks: list[Block] = []
    current_heading: Optional[str] = None
    for element in soup.find_all(
        ["h1", "h2", "h3", "h4", "p", "li", "td", "th"]
    ):
        text = element.get_text(" ", strip=True)
        if not text:
            continue
        if element.name in ("h1", "h2", "h3", "h4"):
            current_heading = text
        blocks.append(Block(text=text, heading=current_heading))
    return ParsedDocument(blocks=blocks, title=title)


def _guess_heading(text: str) -> Optional[str]:
    first_line = text.strip().splitlines()[0] if text.strip() else ""
    if 0 < len(first_line) <= 80 and not first_line.endswith("."):
        return first_line.strip()
    return None


def file_stats(path: str) -> dict:
    st = os.stat(path)
    import datetime

    return {
        "file_size": st.st_size,
        "file_created": datetime.datetime.fromtimestamp(st.st_ctime).isoformat(),
        "file_modified": datetime.datetime.fromtimestamp(st.st_mtime).isoformat(),
    }
